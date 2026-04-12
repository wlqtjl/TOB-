"""
培训材料解析器 — MinerU 2.5 集成

优先使用 MinerU 2.5 进行高精度文档解析 (PDF → 结构化 Markdown)，
支持复杂版面、表格、公式、扫描件 OCR。

降级策略:
  1. MinerU 2.5 (magic_pdf)  — 最优: 版面感知 + 表格 + 公式 + OCR
  2. PyMuPDF (fitz)          — 中等: 纯文本 + 基础结构
  3. python-docx / pptx      — DOCX/PPT 专用解析
"""

from __future__ import annotations

import logging
import tempfile
import shutil
from pathlib import Path
from dataclasses import dataclass, field

logger = logging.getLogger(__name__)

# ── 解析结果数据模型 ─────────────────────────────────────────────────


@dataclass
class ParsedImage:
    """文档中提取的图片"""
    path: str           # 图片临时文件路径
    page: int           # 所在页码 (0-based)
    caption: str = ""   # 图片说明文字


@dataclass
class ParsedTable:
    """文档中提取的表格"""
    html: str           # 表格 HTML
    page: int           # 所在页码
    caption: str = ""   # 表格说明


@dataclass
class ParseResult:
    """文档解析结果"""
    markdown: str                            # 完整 Markdown 文本 (主要输出)
    plain_text: str                          # 纯文本 (兼容旧管道)
    images: list[ParsedImage] = field(default_factory=list)
    tables: list[ParsedTable] = field(default_factory=list)
    page_count: int = 0
    parser_used: str = "unknown"             # mineru / pymupdf / docx / pptx / plaintext
    metadata: dict = field(default_factory=dict)


# ── MinerU 可用性检测 ────────────────────────────────────────────────

_mineru_available: bool | None = None


def is_mineru_available() -> bool:
    """检测 MinerU (magic_pdf) 是否已安装且可用"""
    global _mineru_available
    if _mineru_available is not None:
        return _mineru_available
    try:
        from magic_pdf.data.data_reader_writer import FileBasedDataWriter  # noqa: F401
        from magic_pdf.data.dataset import PymuDocDataset  # noqa: F401
        _mineru_available = True
        logger.info("✅ MinerU (magic_pdf) 可用")
    except ImportError:
        _mineru_available = False
        logger.warning("⚠️ MinerU (magic_pdf) 未安装，将使用 fallback 解析器")
    return _mineru_available


# ── MinerU 2.5 解析 ──────────────────────────────────────────────────


def _parse_pdf_with_mineru(pdf_bytes: bytes, work_dir: Path) -> ParseResult:
    """
    使用 MinerU 2.5 解析 PDF — 高精度版面分析 + 表格 + 公式 + OCR

    MinerU 2.5 内部流程:
      1. 版面分析 (Layout Detection) — 识别文本块/表格/图片/公式区域
      2. 阅读顺序恢复 (Reading Order) — 多栏排序
      3. 内容识别:
         - 文本: 直接提取或 OCR (109+ 语言)
         - 表格: 结构化 HTML 输出
         - 公式: LaTeX 输出
         - 图片: 裁剪保存
      4. Markdown 组装 — 按阅读顺序拼接
    """
    from magic_pdf.data.data_reader_writer import FileBasedDataWriter
    from magic_pdf.data.dataset import PymuDocDataset
    from magic_pdf.model.doc_analyze_by_custom_model import doc_analyze
    import magic_pdf.model as model_config

    # 启用 MinerU 内置模型
    model_config.__use_inside_model__ = True

    # 创建输出目录
    image_dir = work_dir / "images"
    image_dir.mkdir(parents=True, exist_ok=True)
    image_writer = FileBasedDataWriter(str(image_dir))

    # 创建数据集
    dataset = PymuDocDataset(pdf_bytes)

    # 自动判断: 文本 PDF vs 扫描件
    classify_result = dataset.classify()
    logger.info(f"MinerU classify: {classify_result}")

    if classify_result == "ocr":
        # 扫描件 / 图片 PDF — 需要 OCR
        infer_result = dataset.apply(doc_analyze, ocr=True)
        pipe_result = infer_result.pipe_ocr_mode(image_writer)
    else:
        # 文本 PDF — 直接提取
        infer_result = dataset.apply(doc_analyze, ocr=False)
        pipe_result = infer_result.pipe_txt_mode(image_writer)

    # 获取 Markdown 输出
    md_content = pipe_result.get_markdown(str(image_dir))

    # 获取结构化内容列表
    try:
        content_list = pipe_result.get_content_list()
    except Exception:
        content_list = []

    # 提取图片信息
    images: list[ParsedImage] = []
    for img_path in image_dir.glob("*"):
        if img_path.suffix.lower() in (".png", ".jpg", ".jpeg", ".webp"):
            images.append(ParsedImage(
                path=str(img_path),
                page=0,  # MinerU 图片名通常含页码信息
            ))

    # 提取表格
    tables: list[ParsedTable] = []
    for item in content_list:
        if isinstance(item, dict) and item.get("type") == "table":
            tables.append(ParsedTable(
                html=item.get("html", ""),
                page=item.get("page_idx", 0),
            ))

    # 纯文本: 从 Markdown 中去除格式标记
    plain_text = _markdown_to_plain_text(md_content)

    return ParseResult(
        markdown=md_content,
        plain_text=plain_text,
        images=images,
        tables=tables,
        page_count=len(dataset) if hasattr(dataset, "__len__") else 0,
        parser_used="mineru",
        metadata={
            "classify": classify_result,
            "content_blocks": len(content_list),
        },
    )


# ── PyMuPDF Fallback ─────────────────────────────────────────────────


def _parse_pdf_with_pymupdf(pdf_bytes: bytes, work_dir: Path) -> ParseResult:
    """PyMuPDF (fitz) fallback — 基础文本提取"""
    import fitz  # PyMuPDF

    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    pages_text: list[str] = []
    images: list[ParsedImage] = []
    tables: list[ParsedTable] = []

    image_dir = work_dir / "images"
    image_dir.mkdir(parents=True, exist_ok=True)

    for page_idx, page in enumerate(doc):
        # 提取文本 (保留基本格式)
        text = page.get_text("text")
        if text.strip():
            pages_text.append(text.strip())

        # 提取图片
        for img_idx, img_info in enumerate(page.get_images(full=True)):
            try:
                xref = img_info[0]
                base_image = doc.extract_image(xref)
                if base_image and base_image.get("image"):
                    ext = base_image.get("ext", "png")
                    img_path = image_dir / f"page{page_idx}_img{img_idx}.{ext}"
                    img_path.write_bytes(base_image["image"])
                    images.append(ParsedImage(path=str(img_path), page=page_idx))
            except Exception as e:
                logger.debug(f"提取图片失败 (page {page_idx}, img {img_idx}): {e}")

    doc.close()

    plain_text = "\n\n".join(pages_text)

    # 构建简易 Markdown
    md_parts: list[str] = []
    for idx, text in enumerate(pages_text):
        if idx > 0:
            md_parts.append("---")
        md_parts.append(text)

    return ParseResult(
        markdown="\n\n".join(md_parts),
        plain_text=plain_text,
        images=images,
        tables=tables,
        page_count=len(pages_text),
        parser_used="pymupdf",
    )


# ── DOCX 解析 ─────────────────────────────────────────────────────────


def _parse_docx(file_bytes: bytes, work_dir: Path) -> ParseResult:
    """使用 python-docx 解析 Word 文档"""
    import io
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))

    md_parts: list[str] = []
    plain_parts: list[str] = []
    tables: list[ParsedTable] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # 根据样式转换为 Markdown 标题
        style_name = (para.style.name or "").lower()
        if "heading 1" in style_name:
            md_parts.append(f"# {text}")
        elif "heading 2" in style_name:
            md_parts.append(f"## {text}")
        elif "heading 3" in style_name:
            md_parts.append(f"### {text}")
        else:
            md_parts.append(text)
        plain_parts.append(text)

    # 提取表格
    for tbl_idx, table in enumerate(doc.tables):
        rows_html: list[str] = []
        for row_idx, row in enumerate(table.rows):
            cells = [cell.text.strip() for cell in row.cells]
            tag = "th" if row_idx == 0 else "td"
            cells_html = "".join(f"<{tag}>{c}</{tag}>" for c in cells)
            rows_html.append(f"<tr>{cells_html}</tr>")
        html = f"<table>{''.join(rows_html)}</table>"
        tables.append(ParsedTable(html=html, page=0, caption=f"表格 {tbl_idx + 1}"))

        # 添加表格的 Markdown 表示
        if table.rows:
            header = [cell.text.strip() for cell in table.rows[0].cells]
            md_parts.append("| " + " | ".join(header) + " |")
            md_parts.append("| " + " | ".join(["---"] * len(header)) + " |")
            for row in table.rows[1:]:
                cells = [cell.text.strip() for cell in row.cells]
                md_parts.append("| " + " | ".join(cells) + " |")

    return ParseResult(
        markdown="\n\n".join(md_parts),
        plain_text="\n\n".join(plain_parts),
        tables=tables,
        parser_used="docx",
    )


# ── PPT 解析 ──────────────────────────────────────────────────────────


def _parse_pptx(file_bytes: bytes, work_dir: Path) -> ParseResult:
    """使用 python-pptx 解析 PowerPoint"""
    import io
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))

    md_parts: list[str] = []
    plain_parts: list[str] = []
    images: list[ParsedImage] = []

    image_dir = work_dir / "images"
    image_dir.mkdir(parents=True, exist_ok=True)

    for slide_idx, slide in enumerate(prs.slides):
        slide_texts: list[str] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)

            # 提取图片
            if shape.shape_type == 13:  # Picture
                try:
                    image = shape.image
                    ext = image.content_type.split("/")[-1]
                    img_path = image_dir / f"slide{slide_idx}_{shape.shape_id}.{ext}"
                    img_path.write_bytes(image.blob)
                    images.append(ParsedImage(path=str(img_path), page=slide_idx))
                except Exception as e:
                    logger.debug(f"提取PPT图片失败: {e}")

        if slide_texts:
            # 第一行作为幻灯片标题
            md_parts.append(f"## Slide {slide_idx + 1}: {slide_texts[0]}")
            md_parts.extend(slide_texts[1:])
            plain_parts.extend(slide_texts)

    return ParseResult(
        markdown="\n\n".join(md_parts),
        plain_text="\n\n".join(plain_parts),
        images=images,
        page_count=len(prs.slides),
        parser_used="pptx",
    )


# ── 纯文本 ───────────────────────────────────────────────────────────


def _parse_plain_text(file_bytes: bytes) -> ParseResult:
    """纯文本 / Markdown 文件"""
    text = file_bytes.decode("utf-8", errors="replace").strip()
    return ParseResult(
        markdown=text,
        plain_text=text,
        parser_used="plaintext",
    )


# ── 辅助函数 ──────────────────────────────────────────────────────────


def _markdown_to_plain_text(md: str) -> str:
    """将 Markdown 转为纯文本 (去除格式标记)"""
    import re
    text = md
    # 去除图片标记
    text = re.sub(r"!\[([^\]]*)\]\([^)]+\)", r"\1", text)
    # 去除链接, 保留文字
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    # 去除标题标记
    text = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)
    # 去除加粗/斜体
    text = re.sub(r"\*{1,2}([^*]+)\*{1,2}", r"\1", text)
    # 去除行内代码
    text = re.sub(r"`([^`]+)`", r"\1", text)
    # 去除水平线
    text = re.sub(r"^-{3,}$", "", text, flags=re.MULTILINE)
    return text.strip()


def _detect_mime(filename: str, mime: str = "") -> str:
    """从文件名和 MIME 类型推断文件格式"""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if mime == "application/pdf" or ext == "pdf":
        return "pdf"
    if "wordprocessingml" in mime or ext == "docx":
        return "docx"
    if "presentationml" in mime or ext in ("pptx", "ppt"):
        return "pptx"
    if mime == "text/plain" or ext in ("txt", "md", "markdown"):
        return "text"
    return ext


# ── 公共入口 ──────────────────────────────────────────────────────────


async def parse_document(
    file_bytes: bytes,
    filename: str,
    mimetype: str = "",
    work_dir: str | None = None,
) -> ParseResult:
    """
    解析文档 — 统一入口

    根据文件类型自动选择最佳解析器:
      PDF  → MinerU 2.5 (优先) / PyMuPDF (fallback)
      DOCX → python-docx
      PPTX → python-pptx
      TXT  → 直接读取

    Args:
        file_bytes: 文件二进制内容
        filename:   原始文件名
        mimetype:   MIME 类型 (可选)
        work_dir:   工作目录 (可选, 默认自动创建临时目录)

    Returns:
        ParseResult: 结构化解析结果
    """
    file_type = _detect_mime(filename, mimetype)

    # 创建工作目录
    if work_dir:
        _work_dir = Path(work_dir)
        _work_dir.mkdir(parents=True, exist_ok=True)
        auto_cleanup = False
    else:
        _work_dir = Path(tempfile.mkdtemp(prefix="mineru-"))
        auto_cleanup = True

    try:
        if file_type == "pdf":
            return await _parse_pdf(file_bytes, _work_dir)
        elif file_type == "docx":
            return _parse_docx(file_bytes, _work_dir)
        elif file_type == "pptx":
            return _parse_pptx(file_bytes, _work_dir)
        elif file_type == "text":
            return _parse_plain_text(file_bytes)
        else:
            raise ValueError(
                f"不支持的文件类型: {filename} ({mimetype or file_type})。"
                f"支持: PDF, DOCX, PPTX, TXT, MD"
            )
    finally:
        if auto_cleanup:
            shutil.rmtree(_work_dir, ignore_errors=True)


async def _parse_pdf(pdf_bytes: bytes, work_dir: Path) -> ParseResult:
    """PDF 解析 — MinerU 优先, PyMuPDF fallback"""
    if is_mineru_available():
        try:
            logger.info("使用 MinerU 2.5 解析 PDF…")
            return _parse_pdf_with_mineru(pdf_bytes, work_dir)
        except Exception as e:
            logger.warning(f"MinerU 解析失败，切换到 PyMuPDF: {e}")

    try:
        logger.info("使用 PyMuPDF 解析 PDF…")
        return _parse_pdf_with_pymupdf(pdf_bytes, work_dir)
    except Exception as e:
        raise RuntimeError(f"PDF 解析失败: {e}") from e

